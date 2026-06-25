"""PDF、Office 与图片加载器。"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Iterator

import numpy as np
from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document
from PIL import Image

from ..config import OcrConfig
from .ocr import get_ocr, ocr_text_from_result


class PdfFileLoader(BaseLoader):
    def __init__(self, file_path: str, metadata: dict, ocr: OcrConfig) -> None:
        self.file_path = file_path
        self.metadata = metadata
        self.ocr = ocr

    def lazy_load(self) -> Iterator[Document]:
        import fitz

        pdf = fitz.open(self.file_path)
        try:
            ocr_engine = get_ocr(self.ocr.use_cuda) if self.ocr.enabled else None
            for page_index, page in enumerate(pdf, start=1):
                parts = [page.get_text()]
                if ocr_engine is not None:
                    parts.extend(self._ocr_large_images(pdf, page, ocr_engine))
                text = "\n".join(part for part in parts if part.strip())
                if text.strip():
                    yield Document(
                        page_content=text,
                        metadata={**self.metadata, "page": page_index},
                    )
        finally:
            pdf.close()

    def _ocr_large_images(self, pdf: object, page: object, ocr_engine: object) -> list[str]:
        import fitz

        texts: list[str] = []
        for image in page.get_image_info(xrefs=True):
            xref = image.get("xref")
            if not xref:
                continue
            bbox = image["bbox"]
            width_ratio = (bbox[2] - bbox[0]) / page.rect.width
            height_ratio = (bbox[3] - bbox[1]) / page.rect.height
            if width_ratio < self.ocr.pdf_image_width_ratio or height_ratio < self.ocr.pdf_image_height_ratio:
                continue
            pix = fitz.Pixmap(pdf, xref)
            try:
                array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, -1)
                text = ocr_text_from_result(ocr_engine(array))
                if text:
                    texts.append(text)
            finally:
                pix = None
        return texts


class DocxFileLoader(BaseLoader):
    def __init__(self, file_path: str, metadata: dict, ocr: OcrConfig) -> None:
        self.file_path = file_path
        self.metadata = metadata
        self.ocr = ocr

    def lazy_load(self) -> Iterator[Document]:
        from docx import Document as DocxDocument
        from docx import ImagePart

        doc = DocxDocument(self.file_path)
        parts: list[str] = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                parts.append(paragraph.text.strip())
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text = "\n".join(p.text.strip() for p in cell.paragraphs if p.text.strip())
                    if text:
                        parts.append(text)
        if self.ocr.enabled:
            ocr_engine = get_ocr(self.ocr.use_cuda)
            for related in doc.part.related_parts.values():
                if isinstance(related, ImagePart):
                    image = Image.open(BytesIO(related.blob))
                    text = ocr_text_from_result(ocr_engine(np.array(image)))
                    if text:
                        parts.append(text)
        yield Document(page_content="\n".join(parts), metadata=self.metadata.copy())


class PptxFileLoader(BaseLoader):
    def __init__(self, file_path: str, metadata: dict, ocr: OcrConfig) -> None:
        self.file_path = file_path
        self.metadata = metadata
        self.ocr = ocr

    def lazy_load(self) -> Iterator[Document]:
        from pptx import Presentation

        presentation = Presentation(self.file_path)
        ocr_engine = get_ocr(self.ocr.use_cuda) if self.ocr.enabled else None
        for slide_index, slide in enumerate(presentation.slides, start=1):
            parts: list[str] = []
            for shape in sorted(slide.shapes, key=lambda item: (item.top, item.left)):
                parts.extend(_extract_shape_text(shape, ocr_engine))
            text = "\n".join(part for part in parts if part.strip())
            if text:
                yield Document(
                    page_content=text,
                    metadata={**self.metadata, "slide": slide_index},
                )


class ImageFileLoader(BaseLoader):
    def __init__(self, file_path: str, metadata: dict, ocr: OcrConfig) -> None:
        self.file_path = file_path
        self.metadata = metadata
        self.ocr = ocr

    def lazy_load(self) -> Iterator[Document]:
        if not self.ocr.enabled:
            raise ValueError("OCR 已关闭，无法加载图片文件")
        text = ocr_text_from_result(get_ocr(self.ocr.use_cuda)(self.file_path))
        if text.strip():
            yield Document(page_content=text, metadata=self.metadata.copy())


def _extract_shape_text(shape: object, ocr_engine: object | None) -> list[str]:
    parts: list[str] = []
    if getattr(shape, "has_text_frame", False) and shape.text.strip():
        parts.append(shape.text.strip())
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                text = "\n".join(p.text.strip() for p in cell.text_frame.paragraphs if p.text.strip())
                if text:
                    parts.append(text)
    if getattr(shape, "shape_type", None) == 13 and ocr_engine is not None:
        image = Image.open(BytesIO(shape.image.blob))
        text = ocr_text_from_result(ocr_engine(np.array(image)))
        if text:
            parts.append(text)
    if getattr(shape, "shape_type", None) == 6:
        for child in shape.shapes:
            parts.extend(_extract_shape_text(child, ocr_engine))
    return parts


def file_mtime(path: str) -> float:
    return Path(path).stat().st_mtime
